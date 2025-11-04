"""
Générateur de présentations PowerPoint utilisant la charte graphique Infotel 2025
Version exacte conforme à la charte officielle:
- Couleurs PANTONE exactes (653C, 285C, 645C, 654C)
- Polices Segoe UI (Regular, Semilight, Semibold)
- Logo Infotel officiel (PNG)
- PAS de barre rose en footer
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from typing import Dict, List
import os
import shutil

# Import des couleurs officielles Infotel (source unique de vérité)
from services.common import extract_colors_from_template, get_infotel_fonts

# Chemin vers le logo Infotel officiel (remonter de 2 niveaux)
INFOTEL_LOGO_PATH = os.path.join(
    os.path.dirname(__file__),
    "..",
    "..",
    "assets",
    "logo_infotel.png"
)

# Charger les couleurs et polices officielles Infotel (hardcoded, no template needed)
INFOTEL_COLORS = extract_colors_from_template(None)  # Template path not needed anymore
INFOTEL_FONTS = get_infotel_fonts()

def create_powerpoint_from_template(deck_plan: Dict, output_path: str):
    """
    Créer un fichier PowerPoint conforme à la charte graphique Infotel 2025
    
    Args:
        deck_plan: Plan de présentation au format JSON
        output_path: Chemin où sauvegarder le fichier .pptx
    """
    
    # Créer une nouvelle présentation vide (format 16:9)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    colors = INFOTEL_COLORS
    fonts = INFOTEL_FONTS
    
    # Créer chaque slide selon le plan
    for slide_data in deck_plan.get("slides", []):
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "title":
            _create_infotel_title_slide(prs, slide_data, colors, fonts, deck_plan)
        elif slide_type == "section":
            _create_infotel_section_slide(prs, slide_data, colors, fonts)
        elif slide_type == "content" or slide_type == "bullets":
            _create_infotel_content_slide(prs, slide_data, colors, fonts)
        elif slide_type == "comparison":
            _create_infotel_comparison_slide(prs, slide_data, colors, fonts)
        elif slide_type == "conclusion":
            _create_infotel_conclusion_slide(prs, slide_data, colors, fonts)
        else:
            # Fallback
            _create_infotel_content_slide(prs, slide_data, colors, fonts)
    
    # Sauvegarder
    prs.save(output_path)
    print(f"✅ Présentation PowerPoint créée avec charte Infotel 2025: {output_path}")

def _create_infotel_title_slide(prs: Presentation, slide_data: Dict, colors: Dict, fonts: Dict, deck_plan: Dict):
    """Créer une slide de titre style Infotel - Charte 2025"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Bande bleue en haut (couleur primaire Infotel #005091 PANTONE 653C)
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(1.8)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = colors["primary"]
    header.line.fill.background()
    
    # Titre principal (Segoe UI Semibold)
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2),
        Inches(9), Inches(1.5)
    )
    title_frame = title.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", deck_plan.get("title", ""))
    title_p.font.name = fonts["semibold"]  # Segoe UI Semibold
    title_p.font.size = Pt(44)
    title_p.font.color.rgb = colors["primary"]
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.LEFT
    
    # Sous-titre (Segoe UI Regular)
    subtitle_text = slide_data.get("subtitle") or deck_plan.get("subtitle", "")
    if subtitle_text:
        subtitle = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.8),
            Inches(9), Inches(1)
        )
        subtitle_frame = subtitle.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle_text
        subtitle_p.font.name = fonts["regular"]  # Segoe UI Regular
        subtitle_p.font.size = Pt(20)
        subtitle_p.font.color.rgb = colors["text"]
        subtitle_p.alignment = PP_ALIGN.LEFT
    
    # Logo INFOTEL officiel en bas à droite (PNG)
    if os.path.exists(INFOTEL_LOGO_PATH):
        # Insérer le logo PNG (petite taille pour footer)
        slide.shapes.add_picture(
            INFOTEL_LOGO_PATH,
            Inches(8.5), Inches(5.1),  # Position bas à droite
            height=Inches(0.4)  # Hauteur du logo (largeur auto-proportionnelle)
        )
    else:
        # Fallback: texte si logo introuvable
        logo_box = slide.shapes.add_textbox(
            Inches(7.5), Inches(5),
            Inches(2), Inches(0.5)
        )
        logo_frame = logo_box.text_frame
        logo_p = logo_frame.paragraphs[0]
        logo_p.text = "INFOTEL"
        logo_p.font.name = fonts["semibold"]
        logo_p.font.size = Pt(18)
        logo_p.font.color.rgb = colors["primary"]
        logo_p.font.bold = True
        logo_p.alignment = PP_ALIGN.RIGHT

def _create_infotel_section_slide(prs: Presentation, slide_data: Dict, colors: Dict, fonts: Dict):
    """Créer une slide de section style Infotel - Charte 2025"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond bleu primaire (#005091 PANTONE 653C)
    background = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors["primary"]
    background.line.fill.background()
    
    # Titre centré en blanc (Segoe UI Semibold)
    title = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    title_frame = title.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "Section")
    title_p.font.name = fonts["semibold"]  # Segoe UI Semibold
    title_p.font.size = Pt(40)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER

def _create_infotel_content_slide(prs: Presentation, slide_data: Dict, colors: Dict, fonts: Dict):
    """Créer une slide de contenu style Infotel - Charte 2025"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre bleue en haut (plus petite pour slides de contenu)
    header = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(10), Inches(0.6)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = colors["primary"]
    header.line.fill.background()
    
    # Titre de la slide (Segoe UI Semibold)
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8),
        Inches(9), Inches(0.6)
    )
    title_frame = title.text_frame
    title_frame.word_wrap = False
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "")
    title_p.font.name = fonts["semibold"]  # Segoe UI Semibold
    title_p.font.size = Pt(28)
    title_p.font.color.rgb = colors["primary"]
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.LEFT
    
    # Contenu (bullets) - Segoe UI Regular
    bullets = slide_data.get("bullets", [])
    if bullets:
        content = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8),
            Inches(8.5), Inches(3.2)
        )
        content_frame = content.text_frame
        content_frame.word_wrap = True
        
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = bullet_text
            p.font.name = fonts["regular"]  # Segoe UI Regular
            p.font.size = Pt(18)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(8)
            p.space_after = Pt(8)
            p.level = 0
    
    # Numéro de page en bas à droite (Segoe UI Semilight)
    slide_number = len(prs.slides)
    if slide_number > 1:  # Ne pas mettre de numéro sur la première slide
        page_num_box = slide.shapes.add_textbox(
            Inches(9), Inches(5.2),
            Inches(0.8), Inches(0.3)
        )
        page_num_frame = page_num_box.text_frame
        page_num_p = page_num_frame.paragraphs[0]
        page_num_p.text = str(slide_number)
        page_num_p.font.name = fonts["semilight"]  # Segoe UI Semilight
        page_num_p.font.size = Pt(11)
        page_num_p.font.color.rgb = colors["text_light"]
        page_num_p.alignment = PP_ALIGN.RIGHT

def _create_infotel_comparison_slide(prs: Presentation, slide_data: Dict, colors: Dict, fonts: Dict):
    """Créer une slide de comparaison 2 colonnes style Infotel - Charte 2025"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre bleue en haut
    header = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(10), Inches(0.6)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = colors["primary"]
    header.line.fill.background()
    
    # Titre (Segoe UI Semibold)
    title = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8),
        Inches(9), Inches(0.6)
    )
    title_frame = title.text_frame
    title_frame.word_wrap = False
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "")
    title_p.font.name = fonts["semibold"]  # Segoe UI Semibold
    title_p.font.size = Pt(28)
    title_p.font.color.rgb = colors["primary"]
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.LEFT
    
    # Colonne gauche (Segoe UI Regular)
    left_bullets = slide_data.get("left_bullets", [])
    if left_bullets:
        left_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8),
            Inches(4), Inches(3.2)
        )
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        
        for i, bullet_text in enumerate(left_bullets):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            
            p.text = bullet_text
            p.font.name = fonts["regular"]  # Segoe UI Regular
            p.font.size = Pt(16)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(6)
            p.space_after = Pt(6)
    
    # Colonne droite (Segoe UI Regular)
    right_bullets = slide_data.get("right_bullets", [])
    if right_bullets:
        right_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.8),
            Inches(4), Inches(3.2)
        )
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, bullet_text in enumerate(right_bullets):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            
            p.text = bullet_text
            p.font.name = fonts["regular"]  # Segoe UI Regular
            p.font.size = Pt(16)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(6)
            p.space_after = Pt(6)

def _create_infotel_conclusion_slide(prs: Presentation, slide_data: Dict, colors: Dict, fonts: Dict):
    """Créer une slide de conclusion style Infotel - Charte 2025"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond avec dégradé de bleu (simulé avec couleur unie bleu primaire)
    background = slide.shapes.add_shape(
        1,
        Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = colors["primary"]
    background.line.fill.background()
    
    # Titre en blanc (Segoe UI Semibold)
    title = slide.shapes.add_textbox(
        Inches(1), Inches(1.5),
        Inches(8), Inches(1)
    )
    title_frame = title.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "Conclusion")
    title_p.font.name = fonts["semibold"]  # Segoe UI Semibold
    title_p.font.size = Pt(36)
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.bold = True
    title_p.alignment = PP_ALIGN.CENTER
    
    # Points clés en blanc (Segoe UI Regular)
    bullets = slide_data.get("bullets", [])
    if bullets:
        content = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.8),
            Inches(7), Inches(2)
        )
        content_frame = content.text_frame
        content_frame.word_wrap = True
        
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()
            
            p.text = bullet_text
            p.font.name = fonts["regular"]  # Segoe UI Regular
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.space_before = Pt(10)
            p.space_after = Pt(10)
            p.alignment = PP_ALIGN.CENTER
    
    # Logo INFOTEL officiel en bas (PNG - version blanche pour fond bleu)
    if os.path.exists(INFOTEL_LOGO_PATH):
        # Note: Le logo PNG est coloré, mais sur fond bleu il reste lisible
        slide.shapes.add_picture(
            INFOTEL_LOGO_PATH,
            Inches(8.5), Inches(5.1),  # Position bas à droite
            height=Inches(0.4)  # Hauteur du logo
        )
    else:
        # Fallback: texte blanc si logo introuvable
        logo_box = slide.shapes.add_textbox(
            Inches(7.5), Inches(5),
            Inches(2), Inches(0.5)
        )
        logo_frame = logo_box.text_frame
        logo_p = logo_frame.paragraphs[0]
        logo_p.text = "INFOTEL"
        logo_p.font.name = fonts["semibold"]
        logo_p.font.size = Pt(18)
        logo_p.font.color.rgb = RGBColor(255, 255, 255)
        logo_p.font.bold = True
        logo_p.alignment = PP_ALIGN.RIGHT
