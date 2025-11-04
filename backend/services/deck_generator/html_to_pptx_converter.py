"""
Convertisseur HTML → PowerPoint ÉDITABLE
Parse le HTML structuré et reconstruit nativement avec python-pptx
"""

import re
from typing import Dict, List
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from services.common.extract_infotel_colors import INFOTEL_BRAND_COLORS, INFOTEL_FONTS


def parse_html_to_structure(html_content: str) -> Dict:
    """
    Parse le HTML et extrait la structure des slides
    
    Args:
        html_content: HTML complet de la présentation
    
    Returns:
        Dict avec title, subtitle, slides (list)
    """
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Extraire les métadonnées
    title_tag = soup.find('title')
    title = title_tag.get_text() if title_tag else "Présentation"
    
    header = soup.find('div', class_='presentation-header')
    subtitle = ""
    if header:
        h1 = header.find('h1')
        if h1:
            title = h1.get_text(strip=True)
        meta = header.find('div', class_='meta')
        if meta:
            subtitle = meta.get_text(strip=True).split('•')[0].strip()
    
    # Extraire les slides
    slides = []
    slide_divs = soup.find_all('div', class_='slide')
    
    for slide_div in slide_divs:
        slide_type = slide_div.get('data-slide-type', 'content')
        slide_notes = slide_div.get('data-notes', '')
        
        slide_data = {
            'type': slide_type,
            'notes': slide_notes
        }
        
        # Extraire le titre
        title_elem = slide_div.find(['h1', 'h2'], class_='slide-title')
        if title_elem:
            slide_data['title'] = title_elem.get_text(strip=True)
        
        # Extraire le sous-titre
        subtitle_elem = slide_div.find(['p', 'div'], class_='slide-subtitle')
        if subtitle_elem:
            slide_data['subtitle'] = subtitle_elem.get_text(strip=True)
        
        # Extraire les bullets
        bullets = []
        content_div = slide_div.find('div', class_='slide-content')
        if content_div:
            ul = content_div.find('ul')
            if ul:
                for li in ul.find_all('li', recursive=False):
                    bullet_text = li.get_text(strip=True)
                    if bullet_text:
                        bullets.append(bullet_text)
                        
                        # Sous-bullets
                        sub_ul = li.find('ul')
                        if sub_ul:
                            for sub_li in sub_ul.find_all('li'):
                                sub_bullet = sub_li.get_text(strip=True)
                                if sub_bullet:
                                    bullets.append(f"  {sub_bullet}")  # Indentation
        
        if bullets:
            slide_data['bullets'] = bullets
        
        # Pour les slides de comparaison
        if slide_type == 'comparison':
            comparison_container = slide_div.find('div', class_='comparison-container')
            if comparison_container:
                columns = comparison_container.find_all('div', class_='comparison-column')
                bullets = []
                for col in columns:
                    col_title = col.find('h3')
                    if col_title:
                        bullets.append(f"**{col_title.get_text(strip=True)}**")
                    ul = col.find('ul')
                    if ul:
                        for li in ul.find_all('li'):
                            bullets.append(li.get_text(strip=True))
                slide_data['bullets'] = bullets
        
        # Pour les slides de conclusion
        if slide_type == 'conclusion':
            content_div = slide_div.find('div', class_='slide-content')
            if content_div:
                paragraphs = content_div.find_all('p')
                bullets = [p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)]
                if bullets:
                    slide_data['bullets'] = bullets
        
        # Pour les slides d'auteur
        author_elem = slide_div.find(['p'], class_='slide-author')
        if author_elem:
            slide_data['author'] = author_elem.get_text(strip=True)
        
        date_elem = slide_div.find(['p'], class_='slide-date')
        if date_elem:
            slide_data['date'] = date_elem.get_text(strip=True)
        
        slides.append(slide_data)
    
    return {
        'title': title,
        'subtitle': subtitle,
        'slides': slides
    }


def html_to_editable_pptx(html_content: str, output_path: str) -> str:
    """
    Convertit le HTML en PowerPoint ÉDITABLE natif
    Parse le HTML et reconstruit avec python-pptx + template Infotel
    
    Args:
        html_content: HTML complet de la présentation
        output_path: Chemin du fichier .pptx à créer
    
    Returns:
        Chemin du fichier créé
    """
    print("🔄 [HTML→PPTX] Parsing du HTML...")
    structure = parse_html_to_structure(html_content)
    
    print(f"📊 [HTML→PPTX] {len(structure['slides'])} slides détectées")
    
    # Créer la présentation PowerPoint
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9
    prs.slide_height = Inches(5.625)
    
    # Couleurs Infotel
    bleu_nuit = RGBColor(*INFOTEL_BRAND_COLORS['bleu_nuit']['rgb'])
    bleu_ciel = RGBColor(*INFOTEL_BRAND_COLORS['bleu_ciel']['rgb'])
    bleu_profond = RGBColor(*INFOTEL_BRAND_COLORS['bleu_profond']['rgb'])
    blanc = RGBColor(255, 255, 255)
    gris_texte = RGBColor(51, 51, 51)
    
    # Créer chaque slide
    for idx, slide_data in enumerate(structure['slides']):
        print(f"🎨 [HTML→PPTX] Création slide {idx+1}/{len(structure['slides'])} - Type: {slide_data.get('type', 'content')}")
        
        slide_type = slide_data.get('type', 'content')
        
        # Layout vide pour contrôle total
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Fond et éléments communs selon le type
        if slide_type == 'title':
            create_title_slide(slide, slide_data, bleu_nuit, bleu_profond, blanc)
        
        elif slide_type == 'section':
            create_section_slide(slide, slide_data, bleu_ciel, blanc)
        
        elif slide_type == 'conclusion':
            create_conclusion_slide(slide, slide_data, bleu_profond, bleu_nuit, blanc)
        
        elif slide_type == 'comparison':
            create_comparison_slide(slide, slide_data, bleu_nuit, bleu_ciel, gris_texte)
        
        else:  # content (par défaut)
            create_content_slide(slide, slide_data, bleu_nuit, bleu_ciel, gris_texte)
        
        # Ajouter le header gradient
        add_header_gradient(slide, bleu_nuit, bleu_ciel, bleu_profond)
        
        # Ajouter le logo Infotel
        add_infotel_logo(slide)
        
        # Numéro de slide
        add_slide_number(slide, idx + 1, len(structure['slides']), gris_texte)
    
    # Sauvegarder
    print(f"💾 [HTML→PPTX] Sauvegarde vers {output_path}...")
    prs.save(output_path)
    print(f"✅ [HTML→PPTX] PowerPoint créé avec succès!")
    
    return output_path


def create_title_slide(slide, data, bleu_nuit, bleu_profond, blanc):
    """Crée une slide de type 'title'"""
    # Fond bleu dégradé
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bleu_nuit
    
    # Titre principal
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = data.get('title', '')
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = INFOTEL_FONTS['title']
    paragraph.font.size = Pt(54)
    paragraph.font.bold = True
    paragraph.font.color.rgb = blanc
    
    # Sous-titre
    if data.get('subtitle'):
        left = Inches(1.5)
        top = Inches(3.2)
        width = Inches(7)
        height = Inches(0.6)
        
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = subtitle_box.text_frame
        text_frame.text = data.get('subtitle', '')
        text_frame.word_wrap = True
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = INFOTEL_FONTS['caption']
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = blanc
    
    # Date
    if data.get('date'):
        left = Inches(3)
        top = Inches(4.5)
        width = Inches(4)
        height = Inches(0.4)
        
        date_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = date_box.text_frame
        text_frame.text = data.get('date', '')
        
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.name = INFOTEL_FONTS['caption']
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(200, 200, 200)


def create_section_slide(slide, data, bleu_ciel, blanc):
    """Crée une slide de type 'section'"""
    # Fond bleu ciel
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bleu_ciel
    
    # Titre centré
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = data.get('title', '')
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1  # Middle
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = INFOTEL_FONTS['title']
    paragraph.font.size = Pt(48)
    paragraph.font.bold = True
    paragraph.font.color.rgb = blanc


def create_conclusion_slide(slide, data, bleu_profond, bleu_nuit, blanc):
    """Crée une slide de type 'conclusion'"""
    # Fond bleu dégradé
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bleu_profond
    
    # Titre
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = data.get('title', '')
    text_frame.word_wrap = True
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = INFOTEL_FONTS['title']
    paragraph.font.size = Pt(48)
    paragraph.font.bold = True
    paragraph.font.color.rgb = blanc
    
    # Contenu
    if data.get('bullets'):
        left = Inches(1.5)
        top = Inches(2.8)
        width = Inches(7)
        height = Inches(1.5)
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = content_box.text_frame
        text_frame.text = '\n'.join(data['bullets'])
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = INFOTEL_FONTS['body']
            paragraph.font.size = Pt(22)
            paragraph.font.color.rgb = blanc


def create_comparison_slide(slide, data, bleu_nuit, bleu_ciel, gris_texte):
    """Crée une slide de type 'comparison'"""
    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Titre
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = data.get('title', '')
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = INFOTEL_FONTS['title']
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True
    paragraph.font.color.rgb = bleu_nuit
    
    # Deux colonnes
    bullets = data.get('bullets', [])
    mid = len(bullets) // 2
    
    # Colonne gauche
    left_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.3),
        Inches(4.3), Inches(3.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    left_box.line.color.rgb = bleu_ciel
    left_box.line.width = Pt(3)
    
    left_text = left_box.text_frame
    for bullet in bullets[:mid]:
        p = left_text.add_paragraph()
        p.text = bullet
        p.font.name = INFOTEL_FONTS['body']
        p.font.size = Pt(17)
        p.font.color.rgb = gris_texte
        p.space_before = Pt(6)
    
    # Colonne droite
    right_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(5.2), Inches(1.3),
        Inches(4.3), Inches(3.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    right_box.line.color.rgb = bleu_ciel
    right_box.line.width = Pt(3)
    
    right_text = right_box.text_frame
    for bullet in bullets[mid:]:
        p = right_text.add_paragraph()
        p.text = bullet
        p.font.name = INFOTEL_FONTS['body']
        p.font.size = Pt(17)
        p.font.color.rgb = gris_texte
        p.space_before = Pt(6)


def create_content_slide(slide, data, bleu_nuit, bleu_ciel, gris_texte):
    """Crée une slide de type 'content' (standard)"""
    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Titre
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.7)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.text = data.get('title', '')
    
    paragraph = text_frame.paragraphs[0]
    paragraph.font.name = INFOTEL_FONTS['title']
    paragraph.font.size = Pt(36)
    paragraph.font.bold = True
    paragraph.font.color.rgb = bleu_nuit
    
    # Ligne de séparation
    line = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(1.25),
        Inches(9), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = bleu_ciel
    line.line.fill.background()
    
    # Sous-titre (optionnel)
    current_top = Inches(1.4)
    if data.get('subtitle'):
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), current_top, Inches(9), Inches(0.4))
        text_frame = subtitle_box.text_frame
        text_frame.text = data.get('subtitle', '')
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = INFOTEL_FONTS['caption']
        paragraph.font.size = Pt(18)
        paragraph.font.italic = True
        paragraph.font.color.rgb = RGBColor(*INFOTEL_BRAND_COLORS['bleu_profond']['rgb'])
        
        current_top = Inches(1.9)
    
    # Bullets
    if data.get('bullets'):
        left = Inches(0.8)
        width = Inches(8.5)
        height = Inches(3.5)
        
        bullets_box = slide.shapes.add_textbox(left, current_top, width, height)
        text_frame = bullets_box.text_frame
        text_frame.word_wrap = True
        
        for bullet_text in data['bullets']:
            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.font.name = INFOTEL_FONTS['body']
            p.font.size = Pt(20)
            p.font.color.rgb = gris_texte
            p.space_before = Pt(9)
            p.level = 1 if bullet_text.startswith('  ') else 0


def add_header_gradient(slide, bleu_nuit, bleu_ciel, bleu_profond):
    """Ajoute la bande de gradient en haut"""
    header = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.08)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = bleu_nuit
    header.line.fill.background()


def add_infotel_logo(slide):
    """Ajoute le texte 'INFOTEL' en bas à droite"""
    left = Inches(8.5)
    top = Inches(5.2)
    width = Inches(1.2)
    height = Inches(0.3)
    
    logo_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = logo_box.text_frame
    text_frame.text = "INFOTEL"
    
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    paragraph.font.name = INFOTEL_FONTS['title']
    paragraph.font.size = Pt(11)
    paragraph.font.color.rgb = RGBColor(*INFOTEL_BRAND_COLORS['bleu_nuit']['rgb'])
    paragraph.font.bold = True


def add_slide_number(slide, current, total, gris_texte):
    """Ajoute le numéro de slide en bas à gauche"""
    left = Inches(0.5)
    top = Inches(5.2)
    width = Inches(1)
    height = Inches(0.3)
    
    number_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = number_box.text_frame
    text_frame.text = f"{current} / {total}"
    
    paragraph = text_frame.paragraphs[0]
    paragraph.font.name = INFOTEL_FONTS['caption']
    paragraph.font.size = Pt(11)
    paragraph.font.color.rgb = RGBColor(128, 128, 128)

