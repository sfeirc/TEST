"""
Constructeur de présentations PowerPoint professionnelles
Crée des fichiers .pptx à partir de plans JSON
FICHIER BACKUP - Utiliser infotel_template_builder.py pour les présentations en production
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List

# Import des couleurs officielles Infotel (source unique de vérité)
from services.common import extract_colors_from_template, get_infotel_fonts

# Charger les couleurs officielles Infotel
INFOTEL_COLORS = extract_colors_from_template("")
INFOTEL_FONTS = get_infotel_fonts()

# Palettes alternatives (utilisant les couleurs Infotel comme base)
COLOR_SCHEMES = {
    "infotel": INFOTEL_COLORS,  # Charte officielle (par défaut)
    "blue": {
        "primary": RGBColor(0, 74, 152),       # Bleu Infotel
        "secondary": RGBColor(0, 159, 227),    # Bleu clair Infotel
        "accent": RGBColor(236, 0, 140),       # Rose Infotel
        "text": RGBColor(51, 51, 51),
        "background": RGBColor(255, 255, 255)
    },
    "modern": {
        "primary": RGBColor(0, 74, 152),       # Bleu Infotel
        "secondary": RGBColor(242, 242, 242),  # Gris clair
        "accent": RGBColor(236, 0, 140),       # Rose Infotel
        "text": RGBColor(51, 51, 51),
        "background": RGBColor(255, 255, 255)
    },
    "corporate": {
        "primary": RGBColor(0, 32, 96),        # Bleu foncé
        "secondary": RGBColor(0, 159, 227),    # Bleu clair Infotel
        "accent": RGBColor(236, 0, 140),       # Rose Infotel
        "text": RGBColor(51, 51, 51),
        "background": RGBColor(255, 255, 255)
    }
}

def create_powerpoint_deck(deck_plan: Dict, output_path: str):
    """
    Créer un fichier PowerPoint à partir d'un plan de présentation
    
    Args:
        deck_plan: Plan de présentation au format JSON
        output_path: Chemin où sauvegarder le fichier .pptx
    """
    
    # Créer une nouvelle présentation
    prs = Presentation()
    prs.slide_width = Inches(10)  # 16:9 aspect ratio
    prs.slide_height = Inches(5.625)
    
    # Obtenir le schéma de couleurs (par défaut: charte Infotel)
    color_scheme_name = deck_plan.get("color_scheme", "infotel")
    colors = COLOR_SCHEMES.get(color_scheme_name, INFOTEL_COLORS)
    
    # Créer chaque slide
    for slide_data in deck_plan.get("slides", []):
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "title":
            _create_title_slide(prs, slide_data, colors, deck_plan)
        elif slide_type == "section":
            _create_section_slide(prs, slide_data, colors)
        elif slide_type == "content" or slide_type == "bullets":
            _create_content_slide(prs, slide_data, colors)
        elif slide_type == "comparison":
            _create_comparison_slide(prs, slide_data, colors)
        elif slide_type == "conclusion":
            _create_conclusion_slide(prs, slide_data, colors)
        else:
            # Fallback: slide de contenu standard
            _create_content_slide(prs, slide_data, colors)
    
    # Sauvegarder
    prs.save(output_path)
    print(f"✅ Présentation PowerPoint créée: {output_path}")

def _create_title_slide(prs: Presentation, slide_data: Dict, colors: Dict, deck_plan: Dict):
    """Créer une slide de titre"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Rectangle de fond coloré en haut
    header_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(2)
    )
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = colors["primary"]
    header_shape.line.fill.background()
    
    # Titre principal
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(9), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = deck_plan.get("title", slide_data.get("title", "Présentation"))
    title_p.font.name = "Arial"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.LEFT
    
    # Sous-titre
    if deck_plan.get("subtitle"):
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = deck_plan.get("subtitle", "")
        subtitle_p.font.name = "Arial"
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = colors["text"]
        subtitle_p.alignment = PP_ALIGN.LEFT
    
    # Logo / Branding (texte "INFOTEL" en bas à droite)
    logo_box = slide.shapes.add_textbox(
        Inches(7.5), Inches(5),
        Inches(2), Inches(0.5)
    )
    logo_frame = logo_box.text_frame
    logo_p = logo_frame.paragraphs[0]
    logo_p.text = "INFOTEL"
    logo_p.font.name = "Arial"
    logo_p.font.size = Pt(20)
    logo_p.font.bold = True
    logo_p.font.color.rgb = colors["primary"]  # Bleu Infotel
    logo_p.alignment = PP_ALIGN.RIGHT

def _create_section_slide(prs: Presentation, slide_data: Dict, colors: Dict):
    """Créer une slide de séparation de section"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Fond avec bande colorée
    bg_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(1.5),
        Inches(10), Inches(2.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["primary"]
    bg_shape.line.fill.background()
    
    # Titre de section
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "Section")
    title_p.font.name = "Arial"
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER

def _create_content_slide(prs: Presentation, slide_data: Dict, colors: Dict):
    """Créer une slide de contenu avec bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre de titre colorée
    title_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors["primary"]
    title_bar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "Contenu")
    title_p.font.name = "Arial"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Contenu (bullets)
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.2),
        Inches(8.5), Inches(4)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    content_items = slide_data.get("content", [])
    for i, item in enumerate(content_items):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = f"• {item}"
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = colors["text"]
        p.space_after = Pt(12)
        p.level = 0
    
    # Notes
    if slide_data.get("notes"):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.get("notes", "")

def _create_comparison_slide(prs: Presentation, slide_data: Dict, colors: Dict):
    """Créer une slide de comparaison deux colonnes"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Barre de titre
    title_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = colors["primary"]
    title_bar.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "Comparaison")
    title_p.font.name = "Arial"
    title_p.font.size = Pt(28)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Diviser le contenu en deux colonnes
    content_items = slide_data.get("content", [])
    mid_point = len(content_items) // 2
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2),
        Inches(4.5), Inches(4)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(content_items[:mid_point]):
        if i > 0:
            p = left_frame.add_paragraph()
        else:
            p = left_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = colors["text"]
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.2),
        Inches(4), Inches(4)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(content_items[mid_point:]):
        if i > 0:
            p = right_frame.add_paragraph()
        else:
            p = right_frame.paragraphs[0]
        p.text = f"• {item}"
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.color.rgb = colors["text"]

def _create_conclusion_slide(prs: Presentation, slide_data: Dict, colors: Dict):
    """Créer une slide de conclusion"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Rectangle de fond
    bg_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        Inches(10), Inches(5.625)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = colors["secondary"]
    bg_shape.line.fill.background()
    
    # Titre
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1),
        Inches(8), Inches(1)
    )
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = slide_data.get("title", "Conclusion")
    title_p.font.name = "Arial"
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = colors["primary"]
    title_p.alignment = PP_ALIGN.CENTER
    
    # Points de conclusion
    content_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.2),
        Inches(7), Inches(2.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(slide_data.get("content", [])):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = f"✓ {item}"
        p.font.name = "Arial"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        p.space_after = Pt(16)
        p.alignment = PP_ALIGN.LEFT

