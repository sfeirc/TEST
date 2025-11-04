"""
Extracteur de contenu PowerPoint
Lit un fichier .pptx existant et extrait son contenu structuré
"""
from pptx import Presentation
from typing import Dict, List
import os

def extract_content_from_pptx(pptx_path: str) -> Dict:
    """
    Extraire le contenu structuré d'un fichier PowerPoint existant
    
    Args:
        pptx_path: Chemin vers le fichier .pptx à analyser
    
    Returns:
        Dict contenant le contenu structuré:
        {
            "title": "Titre de la présentation",
            "total_slides": 10,
            "slides": [
                {
                    "slide_number": 1,
                    "title": "Titre de la slide",
                    "content": ["Point 1", "Point 2", ...],
                    "notes": "Notes du présentateur"
                },
                ...
            ]
        }
    """
    
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"Fichier PowerPoint introuvable: {pptx_path}")
    
    prs = Presentation(pptx_path)
    
    extracted_data = {
        "title": "",
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_data = {
            "slide_number": slide_idx,
            "title": "",
            "content": [],
            "notes": ""
        }
        
        # Extraire le titre de la slide
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()
            
            # Si c'est la première slide, c'est probablement le titre principal
            if slide_idx == 1 and not extracted_data["title"]:
                extracted_data["title"] = slide_data["title"]
        
        # Extraire le contenu textuel de toutes les formes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                
                # Éviter de dupliquer le titre
                if text and text != slide_data["title"]:
                    slide_data["content"].append(text)
            
            # Extraire les bullets des text frames
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    bullet_text = paragraph.text.strip()
                    if bullet_text and bullet_text not in slide_data["content"]:
                        slide_data["content"].append(bullet_text)
        
        # Extraire les notes du présentateur
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_data["notes"] = notes_slide.notes_text_frame.text.strip()
        
        extracted_data["slides"].append(slide_data)
    
    return extracted_data


def extract_text_summary_from_pptx(pptx_path: str) -> str:
    """
    Extraire tout le contenu textuel d'un PowerPoint en un seul texte
    Utile pour passer à l'IA pour restructuration
    
    Args:
        pptx_path: Chemin vers le fichier .pptx
    
    Returns:
        str: Contenu textuel complet de la présentation
    """
    
    extracted_data = extract_content_from_pptx(pptx_path)
    
    summary_lines = []
    
    # Titre principal
    if extracted_data["title"]:
        summary_lines.append(f"# {extracted_data['title']}\n")
    
    # Contenu de chaque slide
    for slide in extracted_data["slides"]:
        if slide["title"]:
            summary_lines.append(f"\n## Slide {slide['slide_number']}: {slide['title']}")
        else:
            summary_lines.append(f"\n## Slide {slide['slide_number']}")
        
        for content in slide["content"]:
            summary_lines.append(f"- {content}")
        
        if slide["notes"]:
            summary_lines.append(f"\nNotes: {slide['notes']}")
    
    return "\n".join(summary_lines)


if __name__ == "__main__":
    # Test
    test_pptx = "test.pptx"
    if os.path.exists(test_pptx):
        content = extract_content_from_pptx(test_pptx)
        print(f"📊 Titre: {content['title']}")
        print(f"📄 Nombre de slides: {content['total_slides']}")
        print(f"\n📝 Résumé:")
        print(extract_text_summary_from_pptx(test_pptx))