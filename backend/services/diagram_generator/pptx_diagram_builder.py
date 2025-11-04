"""
PowerPoint Diagram Builder
Converts diagram specifications to editable PowerPoint files
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List
import os

# Shape type mapping
SHAPE_TYPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded-rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "cloud": MSO_SHAPE.CLOUD,
    "cylinder": MSO_SHAPE.CAN,
    "process": MSO_SHAPE.FLOWCHART_PROCESS,
    "decision": MSO_SHAPE.FLOWCHART_DECISION,
    "data": MSO_SHAPE.FLOWCHART_DATA,
}

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def create_powerpoint_diagram(diagram_spec: Dict, output_path: str) -> str:
    """
    Create PowerPoint file from diagram specification
    
    Args:
        diagram_spec: Diagram specification from AI
        output_path: Path to save the PowerPoint file
    
    Returns:
        Path to created file
    """
    print(f"🎨 Creating PowerPoint diagram: {diagram_spec.get('title', 'Diagram')}")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Get color scheme
    colors = diagram_spec.get('color_scheme', {})
    primary_color = hex_to_rgb(colors.get('primary', '#0078D4'))
    text_color = hex_to_rgb(colors.get('text', '#323130'))
    
    # Add title
    title = diagram_spec.get('title', 'Diagram')
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(12.33), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    
    # Track created shapes for connections
    shape_map = {}
    
    # Draw containers first (background)
    containers = diagram_spec.get('containers', [])
    for container in containers:
        nodes_in_container = container.get('nodes', [])
        if not nodes_in_container:
            continue
        
        # Calculate bounding box for nodes in container
        node_positions = []
        for node_id in nodes_in_container:
            node = next((n for n in diagram_spec.get('nodes', []) if n['id'] == node_id), None)
            if node:
                pos = node.get('position', {})
                size = node.get('size', {})
                node_positions.append({
                    'x': pos.get('x', 0),
                    'y': pos.get('y', 0),
                    'w': size.get('width', 200),
                    'h': size.get('height', 100)
                })
        
        if node_positions:
            # Calculate container bounds
            min_x = min(p['x'] for p in node_positions) - 20
            min_y = min(p['y'] for p in node_positions) - 20
            max_x = max(p['x'] + p['w'] for p in node_positions) + 20
            max_y = max(p['y'] + p['h'] for p in node_positions) + 20
            
            # Draw container rectangle
            container_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(min_x / 100), Inches((min_y / 100) + 1),
                Inches((max_x - min_x) / 100), Inches((max_y - min_y) / 100)
            )
            container_shape.fill.solid()
            container_shape.fill.fore_color.rgb = hex_to_rgb(container.get('color', '#F3F2F1'))
            container_shape.line.color.rgb = hex_to_rgb('#D2D0CE')
            container_shape.line.width = Pt(1)
            
            # Add container label
            if container.get('label'):
                label_box = slide.shapes.add_textbox(
                    Inches(min_x / 100), Inches((min_y / 100) + 0.8),
                    Inches((max_x - min_x) / 100), Inches(0.4)
                )
                label_frame = label_box.text_frame
                label_frame.text = container.get('label')
                label_para = label_frame.paragraphs[0]
                label_para.font.size = Pt(14)
                label_para.font.bold = True
                label_para.font.color.rgb = hex_to_rgb('#605E5C')
    
    # Draw nodes
    nodes = diagram_spec.get('nodes', [])
    for node in nodes:
        node_id = node.get('id')
        label = node.get('label', 'Node')
        node_type = node.get('type', 'rounded-rectangle')
        position = node.get('position', {'x': 100, 'y': 100})
        size = node.get('size', {'width': 200, 'height': 100})
        node_color = node.get('color', colors.get('primary', '#0078D4'))
        
        # Convert to inches (positions are in arbitrary units, scale them)
        left = Inches(position['x'] / 100)
        top = Inches((position['y'] / 100) + 1)  # Offset for title
        width = Inches(size['width'] / 100)
        height = Inches(size['height'] / 100)
        
        # Create shape
        shape_type = SHAPE_TYPES.get(node_type, MSO_SHAPE.ROUNDED_RECTANGLE)
        shape = slide.shapes.add_shape(
            shape_type,
            left, top, width, height
        )
        
        # Style shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(node_color)
        shape.line.color.rgb = hex_to_rgb('#FFFFFF')
        shape.line.width = Pt(2)
        
        # Add text
        text_frame = shape.text_frame
        text_frame.text = label
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        
        # Format text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Add description if present
        description = node.get('description', '')
        if description:
            desc_para = text_frame.add_paragraph()
            desc_para.text = description
            desc_para.font.size = Pt(10)
            desc_para.font.bold = False
            desc_para.font.color.rgb = RGBColor(255, 255, 255)
            desc_para.alignment = PP_ALIGN.CENTER
        
        # Store shape for connections
        shape_map[node_id] = shape
    
    # Draw connections (arrows)
    connections = diagram_spec.get('connections', [])
    for conn in connections:
        from_id = conn.get('from')
        to_id = conn.get('to')
        
        if from_id in shape_map and to_id in shape_map:
            from_shape = shape_map[from_id]
            to_shape = shape_map[to_id]
            
            # Create connector
            connector = slide.shapes.add_connector(
                1,  # Straight connector
                from_shape.left + from_shape.width // 2,
                from_shape.top + from_shape.height // 2,
                to_shape.left + to_shape.width // 2,
                to_shape.top + to_shape.height // 2
            )
            
            conn_color = conn.get('color', colors.get('text', '#323130'))
            connector.line.color.rgb = hex_to_rgb(conn_color)
            connector.line.width = Pt(2)
            
            # Add arrow
            connector.line.end_arrow_type = 2  # Arrow
    
    # Add annotations
    annotations = diagram_spec.get('annotations', [])
    for annotation in annotations:
        text = annotation.get('text', '')
        position = annotation.get('position', {'x': 500, 'y': 50})
        
        note_box = slide.shapes.add_textbox(
            Inches(position['x'] / 100),
            Inches((position['y'] / 100) + 1),
            Inches(3), Inches(0.5)
        )
        note_frame = note_box.text_frame
        note_frame.text = text
        note_para = note_frame.paragraphs[0]
        note_para.font.size = Pt(12)
        note_para.font.italic = True
        note_para.font.color.rgb = hex_to_rgb('#605E5C')
    
    # Save presentation
    prs.save(output_path)
    print(f"✅ PowerPoint diagram saved: {output_path}")
    
    return output_path

